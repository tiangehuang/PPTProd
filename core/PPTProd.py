# encoding=utf-8
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import openpyxl as xlsx
import yaml
from tqdm import tqdm


def rgb2hex(r, g, b):
    return '{0:02x}{1:02x}{2:02x}'.format(r, g, b)


def zero2one(n)-> float:
    if n > 1.:
        return 1.
    if n < 0.:
        return 0.
    return n


def table_style(shape):
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def _set_cell_border(cell, border_color="#000000", border_width='0'):
    """ Hack function to enable the setting of border width and border color
        - left border
        - right border
        - top border
        - bottom border
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Left Cell Border
    lnL = SubElement(tcPr, 'a:lnL', w='3175', cap='flat', cmpd='sng', algn='ctr')
    lnL_solidFill = SubElement(lnL, 'a:solidFill')
    lnL_srgbClr = SubElement(lnL_solidFill, 'a:srgbClr', val=border_color)
    lnL_prstDash = SubElement(lnL, 'a:prstDash', val='solid')
    lnL_round_ = SubElement(lnL, 'a:round')
    lnL_headEnd = SubElement(lnL, 'a:headEnd', type='none', w='med', len='med')
    lnL_tailEnd = SubElement(lnL, 'a:tailEnd', type='none', w='med', len='med')

    # Right Cell Border
    lnR = SubElement(tcPr, 'a:lnR', w='3175', cap='flat', cmpd='sng', algn='ctr')
    lnR_solidFill = SubElement(lnR, 'a:solidFill')
    lnR_srgbClr = SubElement(lnR_solidFill, 'a:srgbClr', val=border_color)
    lnR_prstDash = SubElement(lnR, 'a:prstDash', val='solid')
    lnR_round_ = SubElement(lnR, 'a:round')
    lnR_headEnd = SubElement(lnR, 'a:headEnd', type='none', w='med', len='med')
    lnR_tailEnd = SubElement(lnR, 'a:tailEnd', type='none', w='med', len='med')

    # Top Cell Border
    lnT = SubElement(tcPr, 'a:lnT', w='3175', cap='flat', cmpd='sng', algn='ctr')
    lnT_solidFill = SubElement(lnT, 'a:solidFill')
    lnT_srgbClr = SubElement(lnT_solidFill, 'a:srgbClr', val=border_color)
    lnT_prstDash = SubElement(lnT, 'a:prstDash', val='solid')
    lnT_round_ = SubElement(lnT, 'a:round')
    lnT_headEnd = SubElement(lnT, 'a:headEnd', type='none', w='med', len='med')
    lnT_tailEnd = SubElement(lnT, 'a:tailEnd', type='none', w='med', len='med')

    # Bottom Cell Border
    lnB = SubElement(tcPr, 'a:lnB', w='3175', cap='flat', cmpd='sng', algn='ctr')
    lnB_solidFill = SubElement(lnB, 'a:solidFill')
    lnB_srgbClr = SubElement(lnB_solidFill, 'a:srgbClr', val=border_color)
    lnB_prstDash = SubElement(lnB, 'a:prstDash', val='solid')
    lnB_round_ = SubElement(lnB, 'a:round')
    lnB_headEnd = SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
    lnB_tailEnd = SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')


class Tree:
    def __init__(self, serial, num, tree_type, r, age):
        self.serial = str(serial)
        self.num = str(num)
        self.tree_type = str(tree_type)
        self.r = str(r)
        self.age = str(age)


class PPTProd:
    def __init__(self, config_path='conf/config.yaml'):
        print('=====init=====')
        self.list_table_path = None
        self.image_prefix_path = None
        self.image_suffix = None
        self.template_path = None
        self.result_path = None
        self.point_horizontal_pos = None
        self.point_vertical_pos = None
        self.title = None
        self.trees = None
        self.title = None
        self.pbar = None
        self.prs = None
        self.left = Inches(0.9)
        self.top = Inches(1.8)
        self.w = None
        self.h = None
        self.space = None
        self.num_in_slide = None
        self.table_font_size = None
        self.table_font_style = None

        self.load_config(config_path)
        self.init()

    def load_config(self, config_path):
        with open(config_path, 'r', encoding='utf-8') as f:
            config = yaml.load(f, yaml.FullLoader)
            ppt_config = config['ppt']
            self.list_table_path = ppt_config['list_table_path']
            self.image_prefix_path = ppt_config['image_prefix_path']
            self.image_suffix = ppt_config['image_suffix']
            self.result_path = ppt_config['result_path']
            self.point_horizontal_pos = zero2one(ppt_config['point_horizontal_pos'])
            self.point_vertical_pos = zero2one(ppt_config['point_vertical_pos'])
            self.ppt_title = ppt_config['title']
            # self.template_path = ppt_config['template_path']
            self.template_path = 'template/template.pptx'
            self.w = Inches(ppt_config['width'])
            self.h = Inches(ppt_config['height'])
            self.space = Inches(ppt_config['space'])
            self.num_in_slide = ppt_config['num_in_slide']
            self.table_font_size = Pt(ppt_config['table_font_size'])
            self.table_font_style = ppt_config['table_font_style']

    def init(self):
        wb = xlsx.load_workbook(self.list_table_path)
        sheet = wb[wb.sheetnames[0]]
        self.title = sheet[2]
        self.trees = []
        for row in sheet[3:sheet.max_row]:
            serial = row[0].value
            num = row[1].value
            tree_type = row[2].value
            r = row[3].value
            age = row[4].value
            if num is None:
                continue
            tree = Tree(serial, num, tree_type, r, age)
            self.trees.append(tree)
        self.pbar = tqdm(total=len(self.trees))
        self.prs = Presentation(self.template_path)

    def new_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        shape = slide.shapes[0]
        shape.text_frame.text = self.ppt_title
        shape.text_frame.paragraphs[0].font.bold = True
        return slide

    def proc(self):
        count = 0
        slide = self.new_slide()
        for i in range(len(self.trees)):
            if count == self.num_in_slide:
                slide = self.new_slide()
                count = 0
            l0 = self.left + (self.w + self.space) * count
            t0 = self.top
            img_path = self.image_prefix_path + '/' + self.trees[i].num + '.' + self.image_suffix
            pic = slide.shapes.add_picture(img_path, l0, t0, self.w, self.h)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(l0+self.w*self.point_horizontal_pos), int(t0+self.h*self.point_vertical_pos), Inches(.2), Inches(.2))
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
            line = shape.line
            line.color.rgb = RGBColor(255, 0, 0)

            shape = slide.shapes.add_table(2, 2, l0, t0 + self.h + Inches(0.15), self.w, Inches(0.5))
            table_style(shape)
            table = shape.table
            # write column headings
            table.cell(0, 0).text = self.title[0].value + ':' + self.trees[i].serial
            table.cell(0, 1).text = self.trees[i].tree_type
            # write body cells
            tmp = str(self.title[3].value).split('\n/')
            table.cell(1, 0).text = tmp[0] + self.trees[i].r + tmp[1]
            tmp = str(self.title[4].value).split('\n/')
            table.cell(1, 1).text = tmp[0] + self.trees[i].age + tmp[1]
            for cell in iter_cells(table):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(198, 217, 241)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.name = self.table_font_style
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.font.size = self.table_font_size
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                _set_cell_border(cell, rgb2hex(198, 217, 241))
            self.pbar.update(1)
            count = count + 1

        self.prs.save(self.result_path)
